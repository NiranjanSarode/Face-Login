from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse, JsonResponse
from django import forms
# from .models import Inventory, Weather, Movie
import csv
import os
import sqlite3
from transformers import AutoTokenizer, AutoModelForCausalLM
from pptx import Presentation
from pptx.util import Inches
from docx import Document
import json
import xml.etree.ElementTree as ET
from .models import QueryHistory


name = ""

# Define InventoryForm class
# class InventoryForm(forms.ModelForm):
#     class Meta:
#         model = Inventory
#         fields = ['product_name', 'quantity_in_stock', 'cost_per_item', 'quantity_sold', 'sales', 'stock_date', 'photos']
#         widgets = {
#             'stock_date': forms.DateInput(attrs={'type': 'date'}),
#         }

# Define QueryForm class
class QueryForm(forms.Form):
    query = forms.CharField(label='Enter your natural language query', max_length=200)

# Load the ML model and tokenizer
tokenizer = AutoTokenizer.from_pretrained("./nsql")
model = AutoModelForCausalLM.from_pretrained("./nsql")


term_mapping = {
    "raining":"rain",
    "no precipitation": "NULL",
    "precipitation": "rain",
    "snowfall": "snow",
    "breezy and little cloudy": "Breezy and Partly Cloudy",
    "breezy and cloudy": "Breezy and Mostly Cloudy",
    "breezy and overcast": "Breezy and Overcast",
    "humid and little cloudy": "Humid and Partly Cloudy",
    "humid and cloudy": "Humid and Mostly Cloudy",
    "windy and little cloudy": "Windy and Partly Cloudy",
    "dry and little cloudy": "Dry and Partly Cloudy",
    "very windy and little cloudy": "Dangerously Windy and Partly Cloudy",
    "little cloudy": "Partly Cloudy",
    "dry and cloudy": "Dry and Mostly Cloudy",
    "windy and cloudy": "Windy and Mostly Cloudy",
    "cloudy": "Mostly Cloudy",
    
}

# List of valid summary terms from the database schema
valid_summaries = [
    'Partly Cloudy', 'Mostly Cloudy', 'Overcast', 'Foggy', 'Breezy and Mostly Cloudy', 'Clear',
    'Breezy and Partly Cloudy', 'Breezy and Overcast', 'Humid and Mostly Cloudy', 'Humid and Partly Cloudy',
    'Windy and Foggy', 'Windy and Overcast', 'Breezy and Foggy', 'Windy and Partly Cloudy', 'Breezy',
    'Dry and Partly Cloudy', 'Windy and Mostly Cloudy', 'Dangerously Windy and Partly Cloudy', 'Dry',
    'Windy', 'Humid and Overcast', 'Light Rain', 'Drizzle', 'Windy and Dry', 'Dry and Mostly Cloudy',
    'Breezy and Dry', 'Rain'
]

# List of valid precip_type terms from the database schema
valid_precip_types = ['rain', 'snow', 'NULL']

# Function to replace terms in user input if they are valid summary or precip_type terms
def replace_terms(text, mapping, valid_summary_terms, valid_precip_terms):
    for key, value in mapping.items():
        if key in text:
            if value in valid_summary_terms:
                text = text.replace(key, f"summary = '{value}'")
            elif value in valid_precip_terms:
                text = text.replace(key, f"precip_type = '{value}'")
    return text


def query_database(db_path, query):
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()
        cursor.execute(query)
        results = cursor.fetchall()
        column_names = [description[0] for description in cursor.description]
        conn.close()
        return [dict(zip(column_names, row)) for row in results]
    except sqlite3.Error as e:
        return str(e)
    
def get_sql_query(question):
    processed_query = replace_terms(question, term_mapping, valid_summaries, valid_precip_types)

    text = f""" {{
    "myapps_weather": {{
        "columns": {{
            "id":"PRIMARY KEY INTEGER",
            "formatted_date": "datetime",
            "summary": "ENUM('Partly Cloudy','Mostly Cloudy','Overcast', 'Foggy','Breezy and Mostly Cloudy','Clear', 'Breezy and Partly Cloudy','Breezy and Overcast','Humid and Mostly Cloudy','Humid and Partly Cloudy','Windy and Foggy','Windy and Overcast','Breezy and Foggy','Windy and Partly Cloudy','Breezy','Dry and Partly Cloudy','Windy and Mostly Cloudy','Dangerously Windy and Partly Cloudy','Dry','Windy','Humid and Overcast','Light Rain','Drizzle','Windy and Dry','Dry and Mostly Cloudy','Breezy and Dry','Rain')",
            "temperature_c": "REAL",
            "precip_type": "ENUM('rain', 'snow', 'NULL')",
            "apparent_temperature_c": "REAL",
            "humidity":"REAL",
            "wind_speed_kmh":"REAL",
            "wind_bearing_degrees":"INTEGER",
            "visibilty_km":"REAL",
            "loud_cover":"REAL",
            "pressure_millibars":"REAL",
            "daily_summary":"TEXT"
        }}
        }}
    }}

    -- Using valid SQLite, answer the following questions for the tables provided above.

    -- Example 1: temperature when precipitation is expected
    SELECT temperature_c FROM myapps_weather WHERE precip_type = 'rain';

    -- Example 2: temperature when it is clear
    SELECT temperature_c FROM myapps_weather WHERE summary = 'Clear';

    -- Example 3: temperature when it is partly cloudy
    SELECT temperature_c FROM myapps_weather WHERE summary = 'Partly Cloudy';



    -- {processed_query}
    """

    # Tokenize input
    input_ids = tokenizer(text, return_tensors="pt").input_ids

    # Generate output
    generated_ids = model.generate(input_ids, max_length=512)
    generated_text=tokenizer.decode(generated_ids[0], skip_special_tokens=True)
    return generated_text

def download_csv(results):
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = 'attachment; filename="results.csv"'
    writer = csv.writer(response)
    if results:
        headers = results[0].keys()
        writer.writerow(headers)
        for row in results:
            writer.writerow(row.values())
    return response

def download_word(results):
    document = Document()
    document.add_heading('Query Results', 0)
    if results:
        headers = results[0].keys()
        table = document.add_table(rows=1, cols=len(headers))
        hdr_cells = table.rows[0].cells
        for i, header in enumerate(headers):
            hdr_cells[i].text = header
        for row in results:
            row_cells = table.add_row().cells
            for i, value in enumerate(row.values()):
                row_cells[i].text = str(value)
    document_path = 'query_results.docx'
    document.save(document_path)
    with open(document_path, 'rb') as docx_file:
        response = HttpResponse(docx_file.read(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        response['Content-Disposition'] = 'attachment; filename=query_results.docx'
    return response

def download_ppt(results):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Query Results"
    if not results:
        return HttpResponse("No results found")
    headers = results[0].keys()
    table = slide.shapes.add_table(rows=len(results) + 1, cols=len(headers), left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(0.8)).table
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    for i, row in enumerate(results):
        for j, value in enumerate(row.values()):
            table.cell(i + 1, j).text = str(value)
    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename="results.pptx"'
    prs.save(response)
    return response

def download_history_xml(request):
    # query_history = request.session.get('query_history', [])
    user_id = request.user.username
    query_history = QueryHistory.objects.filter(user_id=user_id).values('query_history')
    if query_history:
        query_history = query_history[0]['query_history']
    else:
        query_history = []
    root = ET.Element('QueryHistory')
    for item in query_history:
        entry = ET.SubElement(root, 'Entry')
        question = ET.SubElement(entry, 'Question')
        question.text = item['question']
        sql_query = ET.SubElement(entry, 'SQLQuery')
        sql_query.text = item['sql_query']
        results = ET.SubElement(entry, 'Results')
        if isinstance(item['results'], str):
            result_entry = ET.SubElement(results, 'Result')
            result_field = ET.SubElement(result_entry, 'Error')
            result_field.text = item['results']
        else:
            for result in item['results']:
                result_entry = ET.SubElement(results, 'Result')
                for key, value in result.items():
                    result_field = ET.SubElement(result_entry, key)
                    result_field.text = str(value)
    xml_str = ET.tostring(root, encoding='utf-8')
    response = HttpResponse(xml_str, content_type='application/xml')
    response['Content-Disposition'] = 'attachment; filename=query_history.xml'
    return response



def clear_history(request):
    # if 'query_history' in request.session:
        # del request.session['query_history']
    # delete all query history where user_id = user_id
    user_id = request.user.username
    QueryHistory.objects.filter(user_id=user_id).delete()
    name = request.GET.get('name', 'default_name')
    return redirect(f'/myapps/?name={name}') 

from django.db import connection
def query_view(request):
    name = request.GET.get('name', 'default_name')
    results = None
    sql_query = None
    query_history = []
    user_id = request.user.username
    query_history = QueryHistory.objects.filter(user_id=user_id).values('query_history')
    if query_history:
        query_history = query_history[0]['query_history']
    else:
        query_history = []

    if request.method == 'POST':
        if 'fetch_table_data' in request.POST:
            table_name = request.POST.get('table_name')
            cursor = connection.cursor()
            cursor.execute(f"SELECT * FROM {table_name};")
            rows = cursor.fetchall()
            columns = [col[0] for col in cursor.description]
            results = [dict(zip(columns, row)) for row in rows]
            return JsonResponse(results, safe=False)

        form = QueryForm(request.POST)
        if form.is_valid():
            question = form.cleaned_data['query']
            sql_query = get_sql_query(question)
            sql_query = sql_query.split('\n')[-1]
            db_path = os.path.join('Chinook_Sqlite.sqlite')
            results = query_database(db_path, sql_query)
            query_history.insert(0, {'question': question, 'sql_query': sql_query, 'results': results})
            if QueryHistory.objects.filter(user_id=user_id).exists():
                QueryHistory.objects.filter(user_id=user_id).update(query_history=query_history)
            else:
                QueryHistory.objects.create(user_id=user_id, query_history=query_history)
            if 'csv' in request.POST:
                return download_csv(results)
            elif 'word' in request.POST:
                return download_word(results)
            elif 'ppt' in request.POST:
                return download_ppt(results)
            response_data = {
                'form': form,
                'results': results,
                'query': sql_query,
                'query_history': query_history,
                'results_json': json.dumps(results),  # Pass the results as JSON
                'username': name
            }
            return render(request, 'query_view.html', response_data)
    else:
        form = QueryForm()
        if 'load_history' in request.GET:
            index = int(request.GET['load_history'])
            history_item = query_history[index]
            form = QueryForm(initial={'query': history_item['question']})
            sql_query = history_item['sql_query']
            results = history_item['results']

    # Fetch table names from SQLite database
    cursor = connection.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
    table_names = [row[0] for row in cursor.fetchall()]

    return render(request, 'query_view.html', {
        'table_names': table_names,
        'form': form,
        'results': results,
        'query': sql_query,
        'query_history': query_history,
        'username': name
    })


import os
import re
from datetime import datetime
from django.db import connection, OperationalError
from django.shortcuts import render
from chardet.universaldetector import UniversalDetector

def infer_data_type(value):
    # Function to infer data type from a value
    try:
        float(value)
        if '.' in value:
            return 'FLOAT'
        else:
            return 'INTEGER'
    except ValueError:
        # List of date formats to check against
        date_formats = [
            '%Y-%m-%d', '%m/%d/%Y', '%d-%m-%Y',
            '%Y/%m/%d', '%d/%m/%Y', '%Y.%m.%d',
            '%d.%m.%Y', '%m.%d.%Y'
        ]
        for date_format in date_formats:
            try:
                datetime.strptime(value, date_format)
                return 'DATE'
            except ValueError:
                continue
        return 'TEXT'  # Default to VARCHAR

def detect_encoding(file):
    detector = UniversalDetector()
    for line in file:
        detector.feed(line)
        if detector.done:
            break
    detector.close()
    return detector.result['encoding']

def upload_csv_and_create_table(request):
    if request.method == 'POST':
        # Handle file upload
        csv_file = request.FILES.get('csv_file')
        if not csv_file:
            return render(request, 'upload_csv.html', {'error_message': 'No file uploaded.'})

        if not csv_file.name.endswith('.csv'):
            return render(request, 'upload_csv.html', {'error_message': 'File is not CSV type'})

        # Detect encoding
        raw_data = csv_file.read()
        encoding = detect_encoding(raw_data.splitlines(True))

        # Read CSV file with detected encoding
        decoded_file = raw_data.decode(encoding).splitlines()
        reader = csv.reader(decoded_file)

        # Extract table name from CSV file name
        table_name = os.path.splitext(csv_file.name)[0]

        # Read CSV file
        # decoded_file = csv_file.read().decode('utf-8').splitlines()
        # reader = csv.reader(decoded_file)

        # Extract headers (first row) from CSV file
        headers = next(reader)

        # Infer data types for each column based on sample values
        sample_values = next(reader)
        column_types = [infer_data_type(value) for value in sample_values]

        # Create table in database using raw SQL
        # Create table in database using raw SQL
        with connection.cursor() as cursor:
            columns = ', '.join([f'"{header}" {data_type}' for header, data_type in zip(headers, column_types)])
            sql_create_table = f'CREATE TABLE IF NOT EXISTS "{table_name}" ({columns})'
            try:
                cursor.execute(sql_create_table)

                for row in reader:
                    values = ', '.join([f"'{value}'" for value in row])
                    columns_str = ', '.join([f'"{header}"' for header in headers])
                    sql_insert_data = f'INSERT INTO "{table_name}" ({columns_str}) VALUES ({values})'
                    cursor.execute(sql_insert_data)

                # Generate schema information
                schema = []
                distinct_threshold = 15  # Adjust this threshold based on your needs
                for header, data_type in zip(headers, column_types):
                    column_info = {
                        'column_name': header,
                        'data_type': data_type,
                        'additional_details': None  # Placeholder for additional details
                    }
                    if data_type == 'TEXT':
                        cursor.execute(f'SELECT COUNT(DISTINCT "{header}") FROM "{table_name}"')
                        distinct_count = cursor.fetchone()[0]
                        if distinct_count <= distinct_threshold:
                            cursor.execute(f'SELECT DISTINCT "{header}" FROM "{table_name}"')
                            distinct_values = [row[0] for row in cursor.fetchall()]
                            column_info['additional_details'] = {'distinct_values': distinct_values}
                    elif data_type in ['INTEGER', 'FLOAT']:
                        cursor.execute(f'SELECT MIN("{header}"), MAX("{header}") FROM "{table_name}"')
                        min_max = cursor.fetchone()
                        column_info['additional_details'] = {
                            'min_value': min_max[0],
                            'max_value': min_max[1]
                        }

                    schema.append(column_info)

                return render(request, 'upload_csv.html', {
                    'success_message': f"CSV file '{csv_file.name}' uploaded and table '{table_name}' created successfully.",
                    'schema': schema,
                    'table_name': table_name,
                    'edit_mode': True,
                })

            except OperationalError as e:
                error_message = f"Operational error: {e}"
                return render(request, 'upload_csv.html', {'error_message': error_message})

    return render(request, 'upload_csv.html')

def save_edited_schema(request):
    if request.method == 'POST':
        table_name = request.POST.get('table_name')
        new_columns = []

        for key in request.POST.keys():
            if key.startswith('column_name_'):
                column_name = request.POST[key]
                data_type = request.POST.get(f'data_type_{key.split("_")[-1]}')
                new_columns.append((column_name, data_type))

        with connection.cursor() as cursor:
            cursor.execute(f"DROP TABLE IF EXISTS \"{table_name}\"")
            columns = ', '.join([f'"{name}" {data_type}' for name, data_type in new_columns])
            sql_create_table = f"CREATE TABLE IF NOT EXISTS \"{table_name}\" ({columns})"
            cursor.execute(sql_create_table)

        return HttpResponse(f"Updated schema for table '{table_name}' saved successfully.")

    return redirect('upload_csv_and_create_table')